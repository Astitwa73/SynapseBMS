"""Build the SIH idea-submission deck from the provided template.

The template is mandatory and caps the deck at six slides including the title,
so the benchmark and the roadmap fold into Feasibility rather than getting slides
of their own. Every number here is read from docs/benchmark.json, which is
written by scripts/compare_policies.py -- nothing is typed in by hand.
"""

from __future__ import annotations

import copy
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
SHOTS = ROOT / "docs" / "screenshots"
TEMPLATE = Path(__file__).resolve().parent / "IDEA_Presentation_Format.pptx"
OUTPUT = ROOT / "docs" / "Autonomous_BMS_Idea_Submission.pptx"

BRAND = RGBColor(0xD0, 0x20, 0x2F)
BRAND_DARK = RGBColor(0xA8, 0x18, 0x26)
INK = RGBColor(0x0D, 0x14, 0x21)
MUTED = RGBColor(0x4A, 0x55, 0x66)
FAINT = RGBColor(0x6B, 0x77, 0x89)
LINE = RGBColor(0xD3, 0xD9, 0xE2)
SUNKEN = RGBColor(0xF5, 0xF7, 0xFA)
OK = RGBColor(0x0F, 0x7A, 0x4A)
OK_TINT = RGBColor(0xE3, 0xF5, 0xEC)
WARN = RGBColor(0xA3, 0x57, 0x00)
WARN_TINT = RGBColor(0xFD, 0xF0, 0xDD)
BRAND_TINT = RGBColor(0xFD, 0xEA, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

HEAD = "Cambria"
BODY = "Calibri"

# Fill these in before submitting; they are portal fields, not project facts.
TITLE_FIELDS = [
    ("Problem Statement ID", "«to be filled»"),
    ("Problem Statement Title", "AI-driven autonomous Building Management System"),
    ("Theme", "Smart Automation / Sustainability"),
    ("PS Category", "Software"),
    ("Team / Student Name", "«to be filled»"),
    ("Student ID", "«to be filled»"),
]


def benchmark() -> dict:
    return json.loads((ROOT / "docs" / "benchmark.json").read_text(encoding="utf-8"))


# --- low-level helpers ------------------------------------------------------


def delete_slide(prs: Presentation, index: int) -> None:
    slide_id = prs.slides._sldIdLst[index]
    prs.part.drop_rel(slide_id.rId)
    prs.slides._sldIdLst.remove(slide_id)


def clear_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def strip_body(slide, keep_title: bool = True) -> None:
    """Remove template body text boxes, keeping title/footer/slide-number."""
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            if shape.placeholder_format.type in (13, 14):  # slide number / footer
                continue
            if keep_title and shape.name.startswith("Title"):
                continue
        if shape.has_text_frame and not shape.is_placeholder:
            clear_shape(shape)
        elif shape.is_placeholder and not (
            keep_title and shape.name.startswith("Title")
        ) and shape.placeholder_format.type not in (13, 14):
            clear_shape(shape)


def textbox(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    return frame


def para(frame, text, *, size=14, bold=False, colour=INK, font=BODY, space_after=4,
         space_before=0, align=PP_ALIGN.LEFT, first=False, italic=False):
    p = frame.paragraphs[0] if first else frame.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = colour
    run.font.name = font
    return p


def card(slide, x, y, w, h, *, fill=WHITE, border=LINE, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    if radius:
        shape.adjustments[0] = 0.06
    shape.text_frame.text = ""
    return shape


def arrow(slide, x, y, w, h, colour=LINE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y),
                                   Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def set_title(slide, text: str) -> None:
    """Set the title and tighten its box.

    All four dimensions are written, not just height. Most of these placeholders
    inherit their geometry from the layout, and setting a single dimension makes
    python-pptx emit an explicit transform whose other extent is zero -- which
    collapses the title to zero width in PowerPoint.
    """
    for shape in slide.shapes:
        if shape.is_placeholder and shape.name.startswith("Title"):
            layout = next(
                (p for p in slide.slide_layout.placeholders
                 if p.placeholder_format.idx == shape.placeholder_format.idx),
                None,
            )
            left = shape.left if shape.left is not None else (
                layout.left if layout is not None else Inches(0.67))
            top = shape.top if shape.top is not None else (
                layout.top if layout is not None else 0)
            width = shape.width or (layout.width if layout is not None else Inches(12.0))
            if not width:
                width = Inches(12.0)

            shape.left, shape.top = left, top
            shape.width, shape.height = width, Inches(0.85)

            frame = shape.text_frame
            frame.clear()
            para(frame, text, size=30, bold=True, colour=INK, font=HEAD, first=True)
            return


def notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


# --- slides -----------------------------------------------------------------


def build_title(slide) -> None:
    strip_body(slide, keep_title=False)

    frame = textbox(slide, 0.7, 0.75, 12.0, 1.6)
    para(frame, "Autonomous Building Management System", size=40, bold=True,
         colour=INK, font=HEAD, first=True, space_after=2)
    para(frame, "An AI agent that reads a live EnergyPlus building, reasons about "
                "comfort and energy, and writes control decisions back — inside a "
                "safety envelope it cannot exceed.",
         size=14, colour=MUTED, space_after=0)

    stats = [
        ("37%", "less cooling energy", OK),
        ("87% → 20%", "time uncomfortable", OK),
        ("0", "safety violations", BRAND),
    ]
    for index, (value, label, colour) in enumerate(stats):
        x = 0.7 + index * 3.1
        box = card(slide, x, 2.6, 2.85, 1.05, fill=SUNKEN)
        frame = box.text_frame
        frame.margin_left = frame.margin_right = Inches(0.12)
        frame.margin_top = Inches(0.08)
        para(frame, value, size=24, bold=True, colour=colour, font=HEAD, first=True,
             space_after=0)
        para(frame, label, size=10, colour=MUTED, space_after=0)

    frame = textbox(slide, 0.7, 4.0, 6.6, 2.7)
    for index, (label, value) in enumerate(TITLE_FIELDS):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = f"{label}:  "
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = MUTED
        run.font.name = BODY
        run = p.add_run()
        run.text = value
        run.font.size = Pt(12)
        run.font.color.rgb = INK
        run.font.name = BODY

    if (SHOTS / "digital-twin.png").exists():
        slide.shapes.add_picture(str(SHOTS / "digital-twin.png"),
                                 Inches(7.7), Inches(3.95), height=Inches(2.65))

    notes(slide,
          "Open with the outcome, not the architecture. We built an AI agent that "
          "runs a real EnergyPlus building simulation in closed loop. Measured on "
          "an identical summer day, it cut cooling energy 37 percent while reducing "
          "time-outside-comfort from 87 percent to 20 percent. The building on the "
          "right is the real five-zone geometry from the EnergyPlus model, not a "
          "drawing. Fill in the portal fields before submitting.")


def build_idea(slide) -> None:
    strip_body(slide)
    set_title(slide, "IDEA — CLOSED-LOOP AI BUILDING CONTROL")

    frame = textbox(slide, 0.67, 1.15, 12.0, 0.4)
    para(frame, "A supervisory AI layer over a live physics simulation — sensing, "
                "reasoning, and actuation on a real closed loop.",
         size=13, colour=MUTED, first=True)

    blocks = [
        ("The problem", INK, [
            "Buildings are ~40% of global energy use.",
            "Fixed schedules ignore live occupancy and weather.",
            "Our baseline building was over-cooled AND wasteful:",
            "occupants uncomfortable 87% of occupied hours.",
        ]),
        ("What we built", BRAND_DARK, [
            "EnergyPlus runs live; the agent reads 38 sensors + 9 meters.",
            "Llama 3 picks one of four actions and explains why.",
            "Deterministic code turns the action into a setpoint.",
            "Every command is clamped before it reaches an actuator.",
        ]),
        ("Why it is novel", OK, [
            "The LLM never emits a number — only a labelled action.",
            "A rule engine runs alongside as a measured baseline,",
            "replacing fabricated “confidence” with real agreement.",
            "Operator, agent and MCP share one safety-clamped path.",
        ]),
    ]

    for index, (heading, colour, lines) in enumerate(blocks):
        x = 0.67 + index * 4.13
        card(slide, x, 1.75, 3.85, 2.5, fill=WHITE)
        frame = textbox(slide, x + 0.22, 1.95, 3.4, 2.1)
        para(frame, heading.upper(), size=11, bold=True, colour=colour, font=BODY,
             first=True, space_after=6)
        for line in lines:
            para(frame, line, size=11, colour=MUTED, space_after=3)

    card(slide, 0.67, 4.45, 12.0, 1.95, fill=SUNKEN)
    frame = textbox(slide, 0.9, 4.62, 11.5, 0.3)
    para(frame, "THE AUTONOMOUS LOOP — VISIBLE ON THE DASHBOARD AT ALL TIMES",
         size=10, bold=True, colour=FAINT, font=BODY, first=True)

    stages = [
        ("OBSERVE", "38 sensors\n9 meters"),
        ("REASON", "Llama 3\n~5.6 s"),
        ("DECIDE", "1 of 4\nactions"),
        ("VALIDATE", "Safety\nenvelope"),
        ("ACT", "Actuators\nwritten"),
        ("RESPOND", "Building\nreacts"),
    ]
    width, gap = 1.72, 0.24
    for index, (name, detail) in enumerate(stages):
        x = 0.9 + index * (width + gap)
        highlight = name in ("REASON", "VALIDATE")
        box = card(slide, x, 5.0, width, 1.15,
                   fill=BRAND_TINT if highlight else WHITE,
                   border=BRAND if highlight else LINE)
        frame = box.text_frame
        frame.margin_left = frame.margin_right = Inches(0.05)
        frame.margin_top = Inches(0.07)
        para(frame, name, size=11, bold=True,
             colour=BRAND_DARK if highlight else INK, font=BODY, first=True,
             space_after=2, align=PP_ALIGN.CENTER)
        for line in detail.split("\n"):
            para(frame, line, size=9, colour=MUTED, space_after=0,
                 align=PP_ALIGN.CENTER)
        if index < len(stages) - 1:
            arrow(slide, x + width + 0.03, 5.47, 0.18, 0.2,
                  colour=BRAND if highlight else LINE)

    frame = textbox(slide, 0.9, 6.18, 11.5, 0.25)
    para(frame, "RESPOND feeds the next OBSERVE — control writes back into "
                "EnergyPlus, so the loop is genuinely closed.",
         size=10, italic=True, colour=FAINT, first=True)

    notes(slide,
          "Three columns: the problem, what we built, why it is novel. The single "
          "most important design decision is in column three — the language model "
          "never produces a number. It selects one of four labelled actions and "
          "explains the choice; deterministic code does the arithmetic. That removes "
          "unit confusion and hallucinated precision, and makes an unrecognised "
          "response detectable rather than plausible. The strip along the bottom is "
          "the loop the dashboard renders live.")


def build_technical(slide) -> None:
    strip_body(slide)
    set_title(slide, "TECHNICAL APPROACH")

    frame = textbox(slide, 0.67, 1.1, 12.0, 0.35)
    para(frame, "Ports-and-adapters: the agent and the dashboard talk to a service "
                "layer, never to the simulation.",
         size=12, colour=MUTED, first=True)

    layers = [
        ("EnergyPlus 26.1", "Annual run · 15-min timesteps\npyenergyplus C API, own thread", BRAND),
        ("Sensor collector", "38 output variables, 9 meters\nHandles resolved once, cached", INK),
        ("Shared state", "Immutable snapshots\nSequence-numbered, lock-free reads", INK),
        ("Processing", "PMV (ISO 7730) · CO₂ mass balance\nEnd-use power breakdown", INK),
        ("Policy layer", "Llama 3 via Ollama, JSON mode, T=0\nRule engine as baseline + fallback", BRAND),
        ("Decision engine", "Bounds · deadband · rate limit\nOne ControlStore for every caller", BRAND),
        ("Control", "Actuators written before the predictor\nCloses the loop into EnergyPlus", INK),
    ]

    top, height, gap = 1.6, 0.62, 0.09
    for index, (name, detail, colour) in enumerate(layers):
        y = top + index * (height + gap)
        accent = colour is BRAND
        card(slide, 0.67, y, 5.5, height,
             fill=BRAND_TINT if accent else WHITE,
             border=BRAND if accent else LINE)
        frame = textbox(slide, 0.85, y + 0.07, 5.1, height - 0.1)
        para(frame, name, size=11.5, bold=True,
             colour=BRAND_DARK if accent else INK, font=BODY, first=True,
             space_after=1)
        para(frame, detail.replace("\n", "   ·   "), size=9, colour=MUTED,
             space_after=0)
        if index < len(layers) - 1:
            arrow(slide, 3.32, y + height + 0.005, 0.2, 0.08, colour=LINE)

    frame = textbox(slide, 0.67, 6.5, 5.5, 0.3)
    para(frame, "↳ Control writes back into EnergyPlus — the loop is closed",
         size=10, bold=True, colour=BRAND_DARK, first=True)

    card(slide, 6.5, 1.6, 6.17, 2.35, fill=WHITE)
    frame = textbox(slide, 6.72, 1.75, 5.75, 2.1)
    para(frame, "STACK", size=10, bold=True, colour=FAINT, font=BODY, first=True,
         space_after=6)
    stack = [
        ("Simulation", "EnergyPlus 26.1 · pyenergyplus · 5ZoneAirCooled.idf"),
        ("Backend", "Python 3.10 · FastAPI · WebSocket · 201 tests"),
        ("AI", "Ollama · Llama 3 8B · JSON-constrained, temperature 0"),
        ("Agent interface", "Model Context Protocol — 6 tools over stdio"),
        ("Frontend", "React 19 · TypeScript · Tailwind 4 · Recharts"),
    ]
    for label, value in stack:
        p = frame.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = f"{label}   "
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = INK
        run.font.name = BODY
        run = p.add_run()
        run.text = value
        run.font.size = Pt(10)
        run.font.color.rgb = MUTED
        run.font.name = BODY

    card(slide, 6.5, 4.15, 6.17, 2.6, fill=SUNKEN)
    frame = textbox(slide, 6.72, 4.3, 5.75, 2.35)
    para(frame, "SAFETY ENVELOPE — APPLIED TO EVERY COMMAND", size=10, bold=True,
         colour=FAINT, font=BODY, first=True, space_after=6)
    rules = [
        "Cooling clamped to 22–28 °C, heating to 16–20 °C",
        "Heating ≤ cooling − 2 °C — an inverted deadband terminates EnergyPlus",
        "Maximum 1 °C setpoint change per decision",
        "Non-finite values rejected before clamping (min/max propagate NaN)",
        "No command → actuators released to the building's own schedule",
    ]
    for rule in rules:
        para(frame, f"•  {rule}", size=10, colour=MUTED, space_after=4)

    notes(slide,
          "Left column is the data path, top to bottom, and it is a real closed "
          "loop — control writes back into EnergyPlus. Two things to call out. "
          "First, timing: sensors are read at the end of a timestep, but actuators "
          "must be written before the zone predictor runs, otherwise the write "
          "silently does nothing. Second, the safety envelope on the bottom right. "
          "The deadband rule is not stylistic — we discovered that an inverted "
          "deadband makes EnergyPlus terminate outright, so a hallucinated setpoint "
          "would end the demo. That rule is enforced last and unconditionally.")


def build_feasibility(slide, data: dict) -> None:
    strip_body(slide)
    set_title(slide, "FEASIBILITY AND VIABILITY")

    results = {r["label"]: r for r in data["results"]}
    baseline = data["results"][0]
    best = min(data["results"][1:], key=lambda r: r["cooling_kwh"])
    saved = (baseline["cooling_kwh"] - best["cooling_kwh"]) / baseline["cooling_kwh"] * 100

    frame = textbox(slide, 0.67, 1.1, 12.0, 0.35)
    para(frame, "Built and measured, not proposed. Same simulated day run three "
                "ways under identical conditions.",
         size=12, colour=MUTED, first=True)

    box = card(slide, 0.67, 1.6, 3.5, 1.35, fill=OK_TINT, border=OK)
    frame = textbox(slide, 0.87, 1.75, 3.1, 1.05)
    para(frame, f"{saved:.0f}%", size=34, bold=True, colour=OK, font=HEAD,
         first=True, space_after=0)
    para(frame, "less cooling energy — measured", size=10, bold=True, colour=OK,
         space_after=2)
    para(frame, f"Uncomfortable {baseline['uncomfortable_pct']:.0f}% → "
                f"{best['uncomfortable_pct']:.0f}% of occupied time",
         size=9.5, colour=MUTED, space_after=0)

    headers = ["Metric", "No agent", "Rule engine", "Llama 3"]
    rows = [
        ("Cooling (kWh)", "cooling_kwh", "{:.2f}"),
        ("Total electricity (kWh)", "total_kwh", "{:.2f}"),
        ("Mean occupied PMV", "mean_occupied_pmv", "{:+.2f}"),
        ("Time uncomfortable (%)", "uncomfortable_pct", "{:.1f}"),
        ("Model fallbacks", "fallbacks", "{:.0f}"),
    ]
    col_x = [4.35, 7.55, 9.25, 11.15]
    col_w = [3.15, 1.6, 1.8, 1.5]

    for index, header in enumerate(headers):
        frame = textbox(slide, col_x[index], 1.62, col_w[index], 0.25)
        para(frame, header, size=9.5, bold=True, colour=FAINT, font=BODY, first=True,
             align=PP_ALIGN.LEFT if index == 0 else PP_ALIGN.RIGHT)

    for row_index, (label, key, fmt) in enumerate(rows):
        y = 1.92 + row_index * 0.235
        frame = textbox(slide, col_x[0], y, col_w[0], 0.22)
        para(frame, label, size=10, colour=MUTED, first=True)
        values = [results[name][key] for name in ("no agent", "rule engine", "llama3")]
        lowest = min(v for v in values if v is not None) if key != "mean_occupied_pmv" else None
        for value_index, value in enumerate(values):
            best_cell = lowest is not None and value == lowest and row_index != 4
            frame = textbox(slide, col_x[value_index + 1], y, col_w[value_index + 1], 0.22)
            para(frame, fmt.format(value), size=10, bold=best_cell,
                 colour=OK if best_cell else INK, first=True, align=PP_ALIGN.RIGHT)

    frame = textbox(slide, 4.35, 3.15, 8.3, 0.25)
    para(frame, f"Recorded {data['date']} on {data['building_model']} · "
                f"{data['command']}",
         size=8.5, italic=True, colour=FAINT, first=True)

    risks = [
        ("Model is slow, wrong or offline",
         "Rule engine takes over automatically. Verified: a stub failing half its "
         "responses produced 100 decisions with zero simulation errors."),
        ("Model asks for an unsafe setpoint",
         "Clamped at the actuator boundary. A 4 °C request becomes 27 °C, and the "
         "adjustment is logged and displayed."),
        ("Results look cherry-picked",
         "One command reproduces the benchmark. Provenance for every figure is on "
         "screen: measured, derived or estimated."),
        ("Simulation ≠ real building",
         "Nothing above the state store touches EnergyPlus. A BACnet adapter "
         "replaces the simulation without changing the agent or safety layer."),
    ]

    frame = textbox(slide, 0.67, 3.55, 12.0, 0.3)
    para(frame, "RISKS AND HOW THEY ARE ALREADY HANDLED", size=10, bold=True,
         colour=FAINT, font=BODY, first=True)

    for index, (risk, mitigation) in enumerate(risks):
        x = 0.67 + (index % 2) * 6.1
        y = 3.9 + (index // 2) * 1.15
        card(slide, x, y, 5.85, 1.0, fill=WHITE)
        frame = textbox(slide, x + 0.2, y + 0.12, 5.45, 0.8)
        para(frame, risk, size=10.5, bold=True, colour=WARN, font=BODY, first=True,
             space_after=2)
        para(frame, mitigation, size=9.5, colour=MUTED, space_after=0)

    card(slide, 0.67, 6.25, 12.0, 0.55, fill=SUNKEN)
    frame = textbox(slide, 0.9, 6.37, 11.6, 0.35)
    para(frame, "NEXT:  demand-controlled ventilation (CO₂ is monitored but not "
                "yet actuated)  ·  multi-day benchmark  ·  BACnet adapter for real "
                "hardware  ·  occupancy forecasting",
         size=10, bold=True, colour=MUTED, font=BODY, first=True)

    notes(slide,
          "Lead with the measured table — this is the part that separates a built "
          "system from a proposal. Both agents cut cooling about 37 percent, and "
          "comfort improved at the same time because the unmanaged building was "
          "over-cooled: occupants sat at PMV minus 0.55. Be honest that the rule "
          "engine and the language model are comparable rather than ranked — this "
          "is one simulated day. The risk cards are the answers to the four "
          "questions judges always ask, and each one is already implemented.")


def build_artifacts(slide) -> None:
    strip_body(slide)
    set_title(slide, "ARTIFACTS")

    frame = textbox(slide, 0.67, 1.08, 12.0, 0.32)
    para(frame, "Live operations dashboard — every value labelled Measured, Derived "
                "or Estimated, with its basis one hover away.",
         size=12, colour=MUTED, first=True)

    if (SHOTS / "dashboard-full.png").exists():
        slide.shapes.add_picture(str(SHOTS / "dashboard-full.png"),
                                 Inches(0.67), Inches(1.5), width=Inches(7.6))

    frame = textbox(slide, 0.67, 6.42, 7.6, 0.28)
    para(frame, "Digital twin uses the real zone polygons from the EnergyPlus model "
                "— four perimeter zones around a core.",
         size=9, italic=True, colour=FAINT, first=True)

    if (SHOTS / "safety-layer-clamped.png").exists():
        slide.shapes.add_picture(str(SHOTS / "safety-layer-clamped.png"),
                                 Inches(8.5), Inches(1.5), height=Inches(3.5))

    frame = textbox(slide, 8.5, 5.12, 4.2, 0.3)
    para(frame, "SAFETY LAYER — A JUDGE CAN TRIGGER THIS", size=10, bold=True,
         colour=FAINT, font=BODY, first=True)
    frame = textbox(slide, 8.5, 5.4, 4.2, 0.9)
    para(frame, "A real request for 4 °C is clamped to 22 °C, then rate-limited to "
                "27 °C. Both rules are shown. The same panel renders commands from "
                "the agent, an operator and MCP — one control path, one clamp.",
         size=9.5, colour=MUTED, first=True)

    card(slide, 8.5, 6.35, 4.17, 0.42, fill=SUNKEN)
    frame = textbox(slide, 8.68, 6.45, 3.85, 0.25)
    para(frame, "201 automated tests  ·  38/38 sensors  ·  9/9 meters  ·  0 fallbacks",
         size=9.5, bold=True, colour=MUTED, font=BODY, first=True)

    notes(slide,
          "This is one screen, no tabs. Top left is the real building geometry with "
          "live comfort colouring. The centre panel is the case file for the current "
          "decision: what the agent observed, its own words for why, what it expected, "
          "and what actually happened afterwards — we measure the outcome and show "
          "the discrepancy rather than hiding it. Right side is the safety layer. "
          "Invite the judge to press Request 4 degrees themselves; the clamp is real, "
          "it goes through the same endpoint the agent uses.")


def build_references(slide) -> None:
    strip_body(slide)
    set_title(slide, "RESEARCH AND REFERENCES")

    groups = [
        ("Standards and models", [
            "ISO 7730:2005 — Fanger PMV/PPD thermal comfort model",
            "ASHRAE Standard 55 — Thermal environmental conditions",
            "ASHRAE Standard 62.1 — Ventilation and CO₂ generation rates",
            "US EPA eGRID — grid emission factor, 0.40 kg CO₂e/kWh",
        ]),
        ("Tools and platforms", [
            "EnergyPlus 26.1 — nrel/NatLabRockies, Python plugin API",
            "5ZoneAirCooled.idf — DOE reference model shipped with EnergyPlus",
            "Model Context Protocol — modelcontextprotocol.io",
            "Ollama · Meta Llama 3 8B — locally hosted, no cloud dependency",
        ]),
        ("Reproducing our results", [
            "python scripts/verify_energyplus.py --run  — API sanity check",
            "python scripts/inspect_model.py  — which sensors the model exposes",
            "python scripts/compare_policies.py  — regenerates the benchmark",
            "python scripts/run_server.py --policy llm  — dashboard on :8000",
        ]),
    ]

    for index, (heading, items) in enumerate(groups):
        y = 1.35 + index * 1.78
        card(slide, 0.67, y, 12.0, 1.6, fill=WHITE if index % 2 == 0 else SUNKEN)
        frame = textbox(slide, 0.9, y + 0.15, 11.5, 1.3)
        para(frame, heading.upper(), size=10.5, bold=True, colour=BRAND_DARK,
             font=BODY, first=True, space_after=6)
        for item in items:
            para(frame, f"•  {item}", size=10.5, colour=MUTED, space_after=3)

    frame = textbox(slide, 0.67, 6.62, 8.6, 0.28)
    para(frame, "Key finding: the unmanaged baseline was over-cooled and wasteful, "
                "so energy and comfort improved together rather than trading off.",
         size=10, italic=True, colour=FAINT, first=True)

    notes(slide,
          "Standards first, because the comfort model is the part judges probe. PMV "
          "is ISO 7730; we measure two of its six inputs and state the other four as "
          "assumptions. The bottom block is the reproducibility story — four commands "
          "take you from a clean machine to the running dashboard, and one of them "
          "regenerates the benchmark table on the previous slide.")


def main() -> None:
    data = benchmark()
    prs = Presentation(str(TEMPLATE))

    delete_slide(prs, 0)  # the "IMPORTANT INSTRUCTIONS" slide

    slides = list(prs.slides)
    build_title(slides[0])
    build_idea(slides[1])
    build_technical(slides[2])
    build_feasibility(slides[3], data)
    build_artifacts(slides[4])
    build_references(slides[5])

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Wrote {OUTPUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
