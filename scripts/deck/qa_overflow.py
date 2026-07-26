"""Estimate text overflow in a deck using real font metrics.

LibreOffice is not available here to render slides, so this measures each text
frame's wrapped height against its box using the actual Windows font files. It
catches the defect class that matters most on a submission deck -- text spilling
past its container -- without a renderer.
"""

from __future__ import annotations

import sys
from pathlib import Path

from PIL import ImageFont
from pptx import Presentation
from pptx.util import Emu

FONT_DIR = Path("C:/Windows/Fonts")
FONT_FILES = {
    "Calibri": ("calibri.ttf", "calibrib.ttf"),
    "Cambria": ("cambria.ttc", "cambriab.ttf"),
    "Arial": ("arial.ttf", "arialbd.ttf"),
}
# Points to pixels at the size PIL is asked to load, plus typical PowerPoint line
# spacing of ~1.2x the point size.
LINE_FACTOR = 1.22
CACHE: dict = {}


def font_for(name: str, size: float, bold: bool):
    key = (name, round(size, 1), bold)
    if key in CACHE:
        return CACHE[key]
    files = FONT_FILES.get(name, FONT_FILES["Calibri"])
    path = FONT_DIR / files[1 if bold else 0]
    if not path.exists():
        path = FONT_DIR / FONT_FILES["Arial"][1 if bold else 0]
    font = ImageFont.truetype(str(path), int(round(size * 96 / 72)))
    CACHE[key] = font
    return font


def wrapped_lines(text: str, font, max_px: float) -> int:
    if not text.strip():
        return 1
    lines, current = 0, ""
    for word in text.split():
        candidate = f"{current} {word}".strip()
        if font.getlength(candidate) <= max_px or not current:
            current = candidate
        else:
            lines += 1
            current = word
    return lines + (1 if current else 0)


def measure(frame, width_in: float) -> float:
    """Estimated rendered height of a text frame, in inches."""
    inner_px = max(10.0, (width_in - 0.04) * 96)
    total_pt = 0.0
    for paragraph in frame.paragraphs:
        runs = paragraph.runs
        if not runs:
            total_pt += 8
            continue
        size = max((r.font.size.pt if r.font.size else 12) for r in runs)
        bold = any(r.font.bold for r in runs)
        name = next((r.font.name for r in runs if r.font.name), "Calibri")
        text = "".join(r.text for r in runs)
        lines = wrapped_lines(text, font_for(name, size, bold), inner_px)
        total_pt += lines * size * LINE_FACTOR
        total_pt += (paragraph.space_after.pt if paragraph.space_after else 0)
        total_pt += (paragraph.space_before.pt if paragraph.space_before else 0)
    return total_pt / 72


def main(path: str) -> int:
    prs = Presentation(path)
    slide_h = Emu(prs.slide_height).inches
    slide_w = Emu(prs.slide_width).inches
    problems = 0

    for index, slide in enumerate(prs.slides, 1):
        findings = []
        for shape in slide.shapes:
            if not shape.has_text_frame or not shape.text_frame.text.strip():
                continue
            w = Emu(shape.width).inches
            h = Emu(shape.height).inches
            top = Emu(shape.top).inches
            left = Emu(shape.left).inches
            needed = measure(shape.text_frame, w)

            if needed > h + 0.06:
                findings.append(
                    f"OVERFLOW  {shape.name!r} needs {needed:.2f}in in {h:.2f}in "
                    f"box :: {shape.text_frame.text[:52]!r}"
                )
            if top + max(h, needed) > slide_h - 0.05:
                findings.append(
                    f"OFF-SLIDE {shape.name!r} bottom at "
                    f"{top + max(h, needed):.2f}in (slide {slide_h:.2f}in)"
                )
            if left + w > slide_w + 0.02:
                findings.append(
                    f"OFF-RIGHT {shape.name!r} right edge at {left + w:.2f}in"
                )
        if findings:
            problems += len(findings)
            print(f"\nSlide {index}:")
            for finding in findings:
                print(f"  {finding}")

    print(f"\n{problems} potential issue(s)")
    return 1 if problems else 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1] if len(sys.argv) > 1 else "out.pptx"))


def overlaps(path: str) -> int:
    """Report shape pairs that overlap enough to be a visible collision.

    Cards deliberately sit behind their own text, so a text frame fully inside a
    background shape is expected. Only partial overlaps between two text-bearing
    shapes, or text crossing a card boundary, are reported.
    """
    prs = Presentation(path)
    found = 0

    def box(shape):
        left = Emu(shape.left).inches
        top = Emu(shape.top).inches
        return (left, top, left + Emu(shape.width).inches, top + Emu(shape.height).inches)

    for index, slide in enumerate(prs.slides, 1):
        texts = [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
        issues = []
        for i, a in enumerate(texts):
            for b in texts[i + 1:]:
                ax1, ay1, ax2, ay2 = box(a)
                bx1, by1, bx2, by2 = box(b)
                ox = min(ax2, bx2) - max(ax1, bx1)
                oy = min(ay2, by2) - max(ay1, by1)
                if ox <= 0.02 or oy <= 0.02:
                    continue
                contained = (ax1 >= bx1 - 0.01 and ax2 <= bx2 + 0.01
                             and ay1 >= by1 - 0.01 and ay2 <= by2 + 0.01)
                reverse = (bx1 >= ax1 - 0.01 and bx2 <= ax2 + 0.01
                           and by1 >= ay1 - 0.01 and by2 <= ay2 + 0.01)
                if contained or reverse:
                    continue
                issues.append(
                    f"OVERLAP {a.name!r} x {b.name!r} by {ox:.2f}x{oy:.2f}in :: "
                    f"{a.text_frame.text[:26]!r} / {b.text_frame.text[:26]!r}"
                )
        if issues:
            found += len(issues)
            print(f"\nSlide {index}:")
            for issue in issues:
                print(f"  {issue}")
    print(f"\n{found} overlap(s)")
    return found
